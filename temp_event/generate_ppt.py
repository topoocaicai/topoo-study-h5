from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# Color palette - black/white/gray premium
# ============================================================
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GRAY = RGBColor(0x99, 0x99, 0x99)

# Font
FONT_TITLE = "Arial Black"
FONT_BODY = "Arial"
FONT_LIGHT = "Arial"

# ============================================================
# Helper functions
# ============================================================
def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, x, y, w, h, font_size=18, font_name=FONT_BODY, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, margin_left=0, margin_right=0):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    return txBox

def add_shape(slide, shape_type, x, y, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_multi_text(slide, lines, x, y, w, h, font_name=FONT_BODY, color=DARK_GRAY, align=PP_ALIGN.LEFT, line_spacing=1.5, margin_left=0):
    """lines: list of tuples (text, font_size, bold, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, fs, b, c = line_info, 14, False, color
        else:
            text = line_info[0]
            fs = line_info[1] if len(line_info) > 1 else 14
            b = line_info[2] if len(line_info) > 2 else False
            c = line_info[3] if len(line_info) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.name = font_name
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = align
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.line_spacing = Pt(fs * line_spacing)
    return txBox

def add_bullet_list(slide, items, x, y, w, h, font_size=13, color=DARK_GRAY, bullet_char="·", line_spacing=1.8, margin_left=0.15):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, item in enumerate(items):
        if isinstance(item, str):
            text, fs, b, c = item, font_size, False, color
        else:
            text = item[0]
            fs = item[1] if len(item) > 1 else font_size
            b = item[2] if len(item) > 2 else False
            c = item[3] if len(item) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {text}"
        p.font.size = Pt(fs)
        p.font.name = FONT_BODY
        p.font.color.rgb = c
        p.font.bold = b
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = Pt(fs * line_spacing)
    return txBox

def add_number_circle(slide, number, x, y, size=0.45, bg_color=BLACK, text_color=WHITE):
    shape = add_shape(slide, MSO_SHAPE.OVAL, x, y, size, size, fill_color=bg_color)
    shape.text_frame.paragraphs[0].text = str(number)
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    shape.text_frame.paragraphs[0].font.name = FONT_BODY
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = False
    shape.text_frame.margin_left = Pt(0)
    shape.text_frame.margin_right = Pt(0)
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape

def slide_number(slide, num, total=12):
    add_text(slide, f"{num} / {total}", 9.0, 5.25, 0.8, 0.3, font_size=9, color=ACCENT_GRAY, align=PP_ALIGN.RIGHT)

def thin_line(slide, x, y, w, color=LIGHT_GRAY, width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(width))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# Create presentation
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

TOTAL_SLIDES = 12

# ============================================================
# Slide 1: Cover (dark)
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(s1, BLACK)
add_text(s1, "TOPOO", 0.8, 0.8, 3, 0.6, font_size=14, font_name=FONT_BODY, color=ACCENT_GRAY, bold=False)
thin_line(s1, 0.8, 1.35, 2.5, ACCENT_GRAY, 0.3)
add_text(s1, "5月南京", 0.8, 1.7, 8, 0.8, font_size=44, font_name=FONT_TITLE, color=WHITE, bold=True)
add_text(s1, "现象级活动方案", 0.8, 2.45, 8, 0.6, font_size=28, font_name=FONT_LIGHT, color=RGBColor(0xAA, 0xAA, 0xAA), bold=False)
add_text(s1, "从0到百万粉丝 — 全网OPC超级个体", 0.8, 3.3, 6, 0.4, font_size=14, font_name=FONT_BODY, color=ACCENT_GRAY)
add_text(s1, "江苏省文化产业协会 全程支持", 0.8, 4.2, 6, 0.3, font_size=11, font_name=FONT_BODY, color=RGBColor(0x66, 0x66, 0x66))
add_text(s1, "2026.05.10 — 05.17  |  南京", 0.8, 4.55, 6, 0.3, font_size=11, font_name=FONT_BODY, color=RGBColor(0x66, 0x66, 0x66))

# ============================================================
# Slide 2: Core Logic
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s2, WHITE)
add_text(s2, "核心逻辑", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s2, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)

add_text(s2, "这不是一场活动，是一个8天的连续事件。", 0.8, 1.5, 8.4, 0.5, font_size=18, color=DARK_GRAY, bold=True)
add_text(s2, "制造一个正在发生的事件，让人觉得错过了就亏了。", 0.8, 2.0, 8.4, 0.5, font_size=14, color=MED_GRAY)

# Three key points
for i, (title, desc) in enumerate([
    ("6场固定场次", "时间分散在8天内，工作日也是全天\n让不同时间有空的人都有机会来"),
    ("1000人必达", "不是冲刺目标，是必须完成\n只能多不能少"),
    ("统一课程内容", "张明军主讲同一套方法论\n越讲越好，口碑越滚越强"),
]):
    x_base = 0.8 + i * 3.0
    add_shape(s2, MSO_SHAPE.RECTANGLE, x_base, 2.8, 2.7, 2.1, fill_color=OFF_WHITE)
    add_number_circle(s2, i + 1, x_base + 0.15, 2.95, 0.4, BLACK, WHITE)
    add_text(s2, title, x_base + 0.7, 2.93, 1.8, 0.35, font_size=15, color=BLACK, bold=True)
    add_text(s2, desc, x_base + 0.15, 3.45, 2.4, 1.2, font_size=11, color=MED_GRAY)

slide_number(s2, 2, TOTAL_SLIDES)

# ============================================================
# Slide 3: Schedule
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s3, WHITE)
add_text(s3, "场次安排", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s3, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)

schedule = [
    ("第一场", "5月10日", "周日", "200人"),
    ("第二场", "5月13日", "周三", "150人"),
    ("第三场", "5月14日", "周四", "150人"),
    ("第四场", "5月15日", "周五", "150人"),
    ("第五场", "5月16日", "周六", "200人"),
    ("第六场", "5月17日", "周日", "收官"),
]

# Header row
add_shape(s3, MSO_SHAPE.RECTANGLE, 0.8, 1.45, 8.4, 0.45, fill_color=BLACK)
for j, header in enumerate(["场次", "日期", "星期", "目标"]):
    add_text(s3, header, 0.8 + j * 2.1, 1.48, 2.1, 0.4, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, (name, date, week, target) in enumerate(schedule):
    y = 1.9 + i * 0.48
    bg_color = OFF_WHITE if i % 2 == 0 else WHITE
    add_shape(s3, MSO_SHAPE.RECTANGLE, 0.8, y, 8.4, 0.48, fill_color=bg_color)
    vals = [name, date, week, target]
    for j, val in enumerate(vals):
        add_text(s3, val, 0.8 + j * 2.1, y + 0.05, 2.1, 0.38, font_size=13, color=DARK_GRAY if j < 3 else BLACK, bold=(j == 3), align=PP_ALIGN.CENTER)

add_text(s3, "6场合计 1050人，超额完成 1000人目标", 0.8, 4.85, 6, 0.3, font_size=13, color=MED_GRAY, bold=True)

slide_number(s3, 3, TOTAL_SLIDES)

# ============================================================
# Slide 4: Instructor
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s4, BLACK)
add_text(s4, "讲师介绍", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=WHITE, bold=True)
thin_line(s4, 0.8, 1.15, 2.5, ACCENT_GRAY, 0.3)

add_text(s4, "张明军", 0.8, 1.6, 6, 0.7, font_size=36, font_name=FONT_TITLE, color=WHITE, bold=True)
add_text(s4, "百万粉丝科技博主  |  AI内容合伙人", 0.8, 2.3, 8, 0.4, font_size=16, color=ACCENT_GRAY)

add_shape(s4, MSO_SHAPE.RECTANGLE, 0.8, 3.0, 8.4, 0.005, fill_color=RGBColor(0x44, 0x44, 0x44))

highlights = [
    "不到一年从0做到百万粉丝，实战打法验证",
    "带出过大量成功学员，方法论可复制",
    "贯穿AI工具赋能，降低内容创作门槛",
    "课程结合南京OPC服务场景，学完即用",
]
for i, h in enumerate(highlights):
    add_text(s4, f"01 0{i+1}   {h}", 0.8, 3.3 + i * 0.45, 8, 0.4, font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))

slide_number(s4, 4, TOTAL_SLIDES)

# ============================================================
# Slide 5: Course Overview
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s5, WHITE)
add_text(s5, "课程大纲", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s5, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)

add_text(s5, "从0到百万粉丝 — 全网OPC超级个体", 0.8, 1.4, 8.4, 0.35, font_size=15, color=MED_GRAY)

# Morning block
add_shape(s5, MSO_SHAPE.RECTANGLE, 0.8, 2.0, 8.4, 0.35, fill_color=LIGHT_GRAY)
add_text(s5, "上午场  09:00 — 12:00   认知破局 + 底层搭建", 1.0, 2.02, 8, 0.3, font_size=13, color=BLACK, bold=True)

morning = [
    ("模块1", "OPC超级个体底层逻辑", "09:00 - 09:40"),
    ("模块2", "短视频IP定位与账号基础搭建", "09:40 - 10:40"),
    ("模块3", "短视频内容创作底层能力", "10:40 - 12:00"),
]
for i, (mod, title, time) in enumerate(morning):
    y = 2.55 + i * 0.52
    add_shape(s5, MSO_SHAPE.RECTANGLE, 1.0, y, 0.7, 0.35, fill_color=DARK_GRAY)
    add_text(s5, mod, 1.0, y + 0.03, 0.7, 0.3, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s5, title, 1.85, y + 0.03, 4.5, 0.3, font_size=13, color=DARK_GRAY)
    add_text(s5, time, 6.8, y + 0.03, 2.4, 0.3, font_size=11, color=ACCENT_GRAY, align=PP_ALIGN.RIGHT)

# Afternoon block
add_shape(s5, MSO_SHAPE.RECTANGLE, 0.8, 4.15, 8.4, 0.35, fill_color=LIGHT_GRAY)
add_text(s5, "下午场  13:30 — 16:30   实操落地 + 商业闭环", 1.0, 4.17, 8, 0.3, font_size=13, color=BLACK, bold=True)

afternoon = [
    ("模块4", "剪辑六步法与AI工具实操", "13:30 - 14:30"),
    ("模块5", "短视频运营与流量冷启动", "14:30 - 15:30"),
    ("模块6", "商业闭环与超级个体变现路径", "15:30 - 16:30"),
]
for i, (mod, title, time) in enumerate(afternoon):
    y = 4.7 + i * 0.3
    add_text(s5, f"{mod}   {title}   {time}", 1.0, y, 8, 0.25, font_size=11, color=DARK_GRAY)

slide_number(s5, 5, TOTAL_SLIDES)

# ============================================================
# Slide 6: Module 1-3 Detail (Morning)
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s6, WHITE)
add_text(s6, "上午场 · 认知破局 + 底层搭建", 0.8, 0.5, 8, 0.6, font_size=24, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s6, 0.8, 1.1, 8.4, LIGHT_GRAY, 0.5)

# Module 1
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 1.35, 8.4, 1.15, fill_color=OFF_WHITE)
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 1.35, 0.06, 1.15, fill_color=BLACK)
add_text(s6, "模块1  OPC超级个体底层逻辑  (09:00 - 09:40)", 1.1, 1.4, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s6, [
    "个人IP成长路径分享 — 百万粉科技博主的从0到1",
    "O（Opportunity）抓住AI+内容时代红利",
    "P（Personal IP）打造高辨识度个人品牌",
    "C（Commercial）实现内容到商业的闭环变现",
], 1.1, 1.72, 7.8, 0.7, font_size=11, color=MED_GRAY)

# Module 2
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 2.65, 8.4, 1.15, fill_color=WHITE)
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 2.65, 0.06, 1.15, fill_color=DARK_GRAY)
add_text(s6, "模块2  短视频IP定位与账号基础搭建  (09:40 - 10:40)", 1.1, 2.7, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s6, [
    "IP定位公式：身份 + 擅长 + 受众，找到差异化标签",
    "抖音/视频号/小红书流量逻辑差异与平台选择",
    "实操：头像、昵称、简介、背景图的「流量密码」设计",
], 1.1, 3.02, 7.8, 0.7, font_size=11, color=MED_GRAY)

# Module 3
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 3.95, 8.4, 1.3, fill_color=OFF_WHITE)
add_shape(s6, MSO_SHAPE.RECTANGLE, 0.8, 3.95, 0.06, 1.3, fill_color=DARK_GRAY)
add_text(s6, "模块3  短视频内容创作底层能力  (10:40 - 12:00)", 1.1, 4.0, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s6, [
    "内容素材库搭建：高效收集、分类、复用素材",
    "文案框架与AI辅助：通用口播结构 + AI生成高质量文案（附提示词模板）",
    "AI镜头与拍摄基础：手机拍摄技巧、运镜方法、AI数字人",
], 1.1, 4.32, 7.8, 0.85, font_size=11, color=MED_GRAY)

slide_number(s6, 6, TOTAL_SLIDES)

# ============================================================
# Slide 7: Module 4-6 Detail (Afternoon)
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s7, WHITE)
add_text(s7, "下午场 · 实操落地 + 商业闭环", 0.8, 0.5, 8, 0.6, font_size=24, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s7, 0.8, 1.1, 8.4, LIGHT_GRAY, 0.5)

# Module 4
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 1.35, 8.4, 1.15, fill_color=OFF_WHITE)
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 1.35, 0.06, 1.15, fill_color=BLACK)
add_text(s7, "模块4  剪辑六步法与AI工具实操  (13:30 - 14:30)", 1.1, 1.4, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s7, [
    "素材处理 → 字幕 → 文案 → 配乐 → 转场 → 导出，完整六步法",
    "AI剪辑工具应用：私教课同款工具演示",
    "实操环节：15分钟快速剪辑，现场点评纠错",
], 1.1, 1.72, 7.8, 0.7, font_size=11, color=MED_GRAY)

# Module 5
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 2.65, 8.4, 1.15, fill_color=WHITE)
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 2.65, 0.06, 1.15, fill_color=DARK_GRAY)
add_text(s7, "模块5  短视频运营与流量冷启动  (14:30 - 15:30)", 1.1, 2.7, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s7, [
    "新号「鱼塘养号法」：冷启动技巧、推高流量池",
    "对标账号拆解、内容矩阵搭建、数据复盘方法",
    "AI优化GEO与内容分发，提升推荐权重",
], 1.1, 3.02, 7.8, 0.7, font_size=11, color=MED_GRAY)

# Module 6
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 3.95, 8.4, 1.3, fill_color=OFF_WHITE)
add_shape(s7, MSO_SHAPE.RECTANGLE, 0.8, 3.95, 0.06, 1.3, fill_color=DARK_GRAY)
add_text(s7, "模块6  商业闭环与超级个体变现路径  (15:30 - 16:30)", 1.1, 4.0, 7.8, 0.3, font_size=13, color=BLACK, bold=True)
add_bullet_list(s7, [
    "短视频商业变现模式：商单、咨询、知识付费、企业服务",
    "私教同款权益解读：合伙人计划、优先承接商单",
    "终极实操：现场梳理个人IP商业计划书，讲师一对一简短点评",
], 1.1, 4.32, 7.8, 0.85, font_size=11, color=MED_GRAY)

slide_number(s7, 7, TOTAL_SLIDES)

# ============================================================
# Slide 8: Course Highlights
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s8, WHITE)
add_text(s8, "课程亮点", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s8, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)

highlights_data = [
    ("全程落地", "上午学理论、搭框架\n下午练实操、做闭环\n避免「听着激动，回去不动」"),
    ("AI赋能", "贯穿AI文案、AI剪辑\nAI数字人工具\n降低创作门槛"),
    ("商业导向", "结合拓圈TOPOO企业服务\n学员IP直接对接\n本地商业需求"),
    ("延续支持", "课后对接深化内容\n长期辅导与商单机会\n完整服务链路"),
]

for i, (title, desc) in enumerate(highlights_data):
    x = 0.8 + i * 2.3
    add_shape(s8, MSO_SHAPE.RECTANGLE, x, 1.6, 2.05, 3.3, fill_color=OFF_WHITE)
    add_shape(s8, MSO_SHAPE.RECTANGLE, x, 1.6, 2.05, 0.04, fill_color=BLACK)
    add_text(s8, title, x + 0.2, 1.85, 1.65, 0.4, font_size=16, color=BLACK, bold=True)
    add_text(s8, desc, x + 0.2, 2.4, 1.65, 2.2, font_size=12, color=MED_GRAY)

slide_number(s8, 8, TOTAL_SLIDES)

# ============================================================
# Slide 9: Propagation Strategy
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s9, WHITE)
add_text(s9, "传播打法", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s9, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)
add_text(s9, "发起人私域信任 × 张明军百万粉丝IP × 江苏省文化产业协会背书 = 信任感", 0.8, 1.4, 8.4, 0.35, font_size=12, color=MED_GRAY)

layers = [
    ("主引擎", "发起人内容输出", "每天围绕活动发内容\n筹备日记、话题抛出、幕后花絮\n倒计时、结束后复盘"),
    ("副引擎", "张明军配合", "百万粉丝账号发预告\n学员群/私域群通知\n联合出镜、观点碰撞"),
    ("核武器", "私域一对一激活", "A类强相关直接私发\nB类可能相关发短版\nC类关系户帮忙转发"),
    ("自增长", "裂变机制", "参与者邀请朋友\n口碑驱动不用利益驱动\n每场结束都是下一场预热"),
]

for i, (tag, title, desc) in enumerate(layers):
    x = 0.8 + i * 2.3
    add_shape(s9, MSO_SHAPE.RECTANGLE, x, 2.1, 2.05, 2.8, fill_color=OFF_WHITE)
    add_shape(s9, MSO_SHAPE.RECTANGLE, x, 2.1, 2.05, 0.04, fill_color=DARK_GRAY)
    add_text(s9, tag, x + 0.15, 2.3, 1.75, 0.25, font_size=10, color=WHITE)
    add_shape(s9, MSO_SHAPE.RECTANGLE, x + 0.15, 2.28, 0.55, 0.25, fill_color=DARK_GRAY)
    add_text(s9, tag, x + 0.15, 2.3, 0.55, 0.22, font_size=10, color=WHITE, bold=True)
    add_text(s9, title, x + 0.15, 2.65, 1.75, 0.35, font_size=15, color=BLACK, bold=True)
    add_text(s9, desc, x + 0.15, 3.05, 1.75, 1.6, font_size=10, color=MED_GRAY)

slide_number(s9, 9, TOTAL_SLIDES)

# ============================================================
# Slide 10: Pricing
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s10, BLACK)
add_text(s10, "定价策略", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=WHITE, bold=True)
thin_line(s10, 0.8, 1.15, 2.5, ACCENT_GRAY, 0.3)

# Two price cards
# Offline
add_shape(s10, MSO_SHAPE.RECTANGLE, 1.2, 1.6, 3.5, 3.0, fill_color=RGBColor(0x22, 0x22, 0x22))
add_text(s10, "线下票", 1.2, 1.8, 3.5, 0.4, font_size=14, color=ACCENT_GRAY, align=PP_ALIGN.CENTER)
add_text(s10, "¥199", 1.2, 2.3, 3.5, 0.8, font_size=48, font_name=FONT_TITLE, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s10, "统一价  不乱价", 1.2, 3.2, 3.5, 0.35, font_size=12, color=ACCENT_GRAY, align=PP_ALIGN.CENTER)
add_text(s10, "6场任意选择  全天沉浸式授课", 1.2, 3.7, 3.5, 0.35, font_size=11, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# Online
add_shape(s10, MSO_SHAPE.RECTANGLE, 5.3, 1.6, 3.5, 3.0, fill_color=RGBColor(0x22, 0x22, 0x22))
add_text(s10, "线上票", 5.3, 1.8, 3.5, 0.4, font_size=14, color=ACCENT_GRAY, align=PP_ALIGN.CENTER)
add_text(s10, "¥99", 5.3, 2.3, 3.5, 0.8, font_size=48, font_name=FONT_TITLE, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s10, "统一价  不乱价", 5.3, 3.2, 3.5, 0.35, font_size=12, color=ACCENT_GRAY, align=PP_ALIGN.CENTER)
add_text(s10, "直播 + 录播回放", 5.3, 3.7, 3.5, 0.35, font_size=11, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# Revenue
add_text(s10, "线下1000人 x ¥199 + 线上300人 x ¥99 = 约 22-23万", 0.8, 4.85, 8.4, 0.3, font_size=12, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

slide_number(s10, 10, TOTAL_SLIDES)

# ============================================================
# Slide 11: Timeline
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s11, WHITE)
add_text(s11, "关键时间轴", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=BLACK, bold=True)
thin_line(s11, 0.8, 1.15, 8.4, LIGHT_GRAY, 0.5)

timeline = [
    ("4.26 - 4.27", "确定6场时间、落实场地"),
    ("4.28", "活动行页面上线，发第一条内容种草"),
    ("4.29 - 4.30", "私域一对一激活开始，张明军发预告"),
    ("5.1 - 5.4", "内容密集输出，裂变开始，目标300人"),
    ("5.5 - 5.9", "最后冲刺，倒计时内容，目标700人"),
    ("5.10（周日）", "第一场 · 首场启动"),
    ("5.10 - 5.17", "边办边传播，每场结束是下一场预热"),
    ("5.17（周日）", "收官场"),
]

for i, (date, action) in enumerate(timeline):
    y = 1.5 + i * 0.48
    bg = OFF_WHITE if i % 2 == 0 else WHITE
    add_shape(s11, MSO_SHAPE.RECTANGLE, 0.8, y, 8.4, 0.48, fill_color=bg)
    add_shape(s11, MSO_SHAPE.RECTANGLE, 0.8, y, 0.04, 0.48, fill_color=BLACK)
    add_text(s11, date, 1.05, y + 0.05, 1.8, 0.35, font_size=11, color=BLACK, bold=True)
    add_text(s11, action, 3.0, y + 0.05, 6.0, 0.35, font_size=12, color=DARK_GRAY)

slide_number(s11, 11, TOTAL_SLIDES)

# ============================================================
# Slide 12: Success Criteria & CTA
# ============================================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s12, BLACK)
add_text(s12, "目标与行动", 0.8, 0.5, 8, 0.6, font_size=28, font_name=FONT_TITLE, color=WHITE, bold=True)
thin_line(s12, 0.8, 1.15, 2.5, ACCENT_GRAY, 0.3)

targets = [
    ("线下累计人数", "1000人（必须完成）"),
    ("活动群规模", "1500人+"),
    ("自发传播", "50条+"),
    ("潜在创造者识别", "100人+"),
    ("拓圈Seeker新增", "50人+"),
]

for i, (label, value) in enumerate(targets):
    y = 1.5 + i * 0.5
    add_text(s12, label, 0.8, y, 4.5, 0.4, font_size=13, color=ACCENT_GRAY)
    add_text(s12, value, 5.5, y, 3.7, 0.4, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

# Bottom CTA
add_shape(s12, MSO_SHAPE.RECTANGLE, 0.8, 4.2, 8.4, 0.5, fill_color=WHITE)
add_text(s12, "这不是活动方案，是拓圈城市落地的标准作战手册", 0.8, 4.23, 8.4, 0.45, font_size=15, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

add_text(s12, "TOPOO  ·  拓圈  ·  让小而美被世界看见", 0.8, 4.95, 8.4, 0.3, font_size=11, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

slide_number(s12, 12, TOTAL_SLIDES)

# ============================================================
# Save
# ============================================================
output_path = r"C:\Users\Administrator\WorkBuddy\Claw\拓圈·5月南京现象级活动方案.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
