# -*- coding: utf-8 -*-
"""
Day 1 视频PPT —— 竖版 9:16
黑白灰高级风格，PPT翻页素材
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 颜色系统 ───
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_BG = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT_RED = RGBColor(0xFF, 0x44, 0x44)
DIM_WHITE = RGBColor(0x99, 0x99, 0x99)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF5)

# 9:16 竖版
SLIDE_W = Inches(9)
SLIDE_H = Inches(16)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── 工具函数 ───
def dark():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    return slide

def light():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE
    return slide

def txt(slide, x, y, w, h, text, sz=36, c=WHITE, b=False, a=PP_ALIGN.CENTER, sp=1.3):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = "Microsoft YaHei"
    p.alignment = a
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if sp != 1.0:
        p.line_spacing = Pt(sz * sp)
    return tb

def multi(slide, x, y, w, h, lines, sz=32, c=WHITE, b=False, a=PP_ALIGN.LEFT, sp=1.6):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            pc = item[1] if len(item) > 1 else c
            pb = item[2] if len(item) > 2 else b
        else:
            p.text = item
            pc, pb = c, b
        p.font.size = Pt(sz)
        p.font.color.rgb = pc
        p.font.bold = pb
        p.font.name = "Microsoft YaHei"
        p.alignment = a
        p.space_after = Pt(2)
        p.line_spacing = Pt(sz * sp)
    return tb

def line(slide, x, y, w, c=MID_GRAY, t=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(t))
    s.fill.solid()
    s.fill.fore_color.rgb = c
    s.line.fill.background()

# ═══════════════════════════════════════════════════
# S1: 钩子
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 3.5, 8, 1.2, "我要做一件", 36, LIGHT_GRAY)
txt(s, 0.5, 4.5, 8, 1.5, "有可能会被打脸的事", 44, WHITE, True)

# ═══════════════════════════════════════════════════
# S2: 自我介绍
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 1, 4.0, 7, 4, [
    ("大家好，我是苏菜菜", WHITE, True),
    ("", WHITE, False),
    ("拓圈Topoo发起者", LIGHT_GRAY, False),
], sz=34, a=PP_ALIGN.CENTER, sp=1.8)

# ═══════════════════════════════════════════════════
# S3: 拓圈不重要 + i人
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 0.8, 3.5, 7.4, 7, [
    ("拓圈是什么不重要", LIGHT_GRAY, False),
    ("以后有时间再讲", DIM_WHITE, False),
    ("", WHITE, False),
    ("真人暂时不出镜", LIGHT_GRAY, False),
    ("是个i人", DIM_WHITE, False),
    ("也等有时间再说", DIM_WHITE, False),
], sz=28, a=PP_ALIGN.CENTER, sp=1.7)

# ═══════════════════════════════════════════════════
# S4: 日期 + 目标
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 2.5, 8, 1.0, "2026年4月26日", 28, MID_GRAY)
line(s, 2, 3.5, 5, MID_GRAY, 1)
multi(s, 0.8, 4.0, 7.4, 7, [
    ("从今天到5月10号", WHITE, False),
    ("14天时间", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("我要在南京找到", WHITE, False),
    ("1000 个人到场", WHITE, True),
    ("", WHITE, False),
    ("做一场关于流量的线下活动", LIGHT_GRAY, False),
], sz=30, a=PP_ALIGN.CENTER, sp=1.6)

# ═══════════════════════════════════════════════════
# S5: 真人坐在椅子上 + 199块
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 0.8, 4.0, 7.4, 7, [
    ("不是线上看直播", WHITE, True),
    ("", WHITE, False),
    ("是真人坐在椅子上", WHITE, True),
    ("待一整天", WHITE, True),
    ("", WHITE, False),
    ("199块", LIGHT_GRAY, False),
], sz=34, a=PP_ALIGN.CENTER, sp=1.6)

# ── 穿插①：空场地实拍 ──

# ═══════════════════════════════════════════════════
# S6: 南京940万
# ═══════════════════════════════════════════════════
s = light()
txt(s, 0.5, 4.0, 8, 1.0, "南京常住人口", 26, MID_GRAY)
txt(s, 0.5, 5.0, 8, 1.5, "940万", 72, RGBColor(0x33, 0x33, 0x33), True)
txt(s, 0.5, 7.0, 8, 1.0, "1000人算什么？", 28, MID_GRAY)

# ═══════════════════════════════════════════════════
# S7: 每一步都在掉人
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 1.5, 8, 0.8, "但做过线下活动的人都知道", 22, LIGHT_GRAY)
line(s, 1.5, 2.5, 6, MID_GRAY, 1)
multi(s, 1.2, 3.0, 6.6, 11, [
    ("让一个人——", WHITE, False),
    ("", WHITE, False),
    ("看到信息", LIGHT_GRAY, False),
    ("↓", MID_GRAY, False),
    ("决定参与", LIGHT_GRAY, False),
    ("↓", MID_GRAY, False),
    ("早上爬起来出门", LIGHT_GRAY, False),
    ("↓", MID_GRAY, False),
    ("坐在椅子上待到下午", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("每一步，都在掉人。", ACCENT_RED, True),
], sz=24, a=PP_ALIGN.CENTER, sp=1.4)

# ── 穿插②：空场地实拍 ──

# ═══════════════════════════════════════════════════
# S8: 零
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 4.5, 8, 1.2, "现在报名人数", 30, LIGHT_GRAY)
txt(s, 0.5, 5.8, 8, 3.0, "零", 120, WHITE, True)

# ── 穿插③：手机录屏/手部特写 ──

# ═══════════════════════════════════════════════════
# S9: 拖延症 + PPT翻页
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 0.8, 3.5, 7.4, 8, [
    ("整个4月份都在盘内容", WHITE, False),
    ("打磨细节", WHITE, False),
    ("", WHITE, False),
    ("直到现在才把素材做好", LIGHT_GRAY, False),
    ("可能有拖延症", DIM_WHITE, False),
    ("", WHITE, False),
    ("也没有时间做视频了", LIGHT_GRAY, False),
    ("PPT翻页吧，我熟", WHITE, True),
], sz=28, a=PP_ALIGN.CENTER, sp=1.7)

# ═══════════════════════════════════════════════════
# S10: 每天更新
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 1, 5.0, 7, 5, [
    ("从今天开始", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("我每天更新", WHITE, True),
    ("真实的信息和数据", WHITE, True),
], sz=32, a=PP_ALIGN.CENTER, sp=1.8)

# ═══════════════════════════════════════════════════
# S11: 做到了算本事
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 1, 4.0, 7, 6, [
    ("做到了算本事", WHITE, True),
    ("", WHITE, False),
    ("没做到", LIGHT_GRAY, False),
    ("这个系列就是大家", LIGHT_GRAY, False),
    ("最好的复盘素材", LIGHT_GRAY, True),
], sz=30, a=PP_ALIGN.CENTER, sp=1.7)

# ═══════════════════════════════════════════════════
# S12: CTA
# ═══════════════════════════════════════════════════
s = dark()
multi(s, 0.8, 4.0, 7.4, 6, [
    ("想看我翻车的", LIGHT_GRAY, False),
    ("点个关注", WHITE, True),
    ("", WHITE, False),
    ("想在南京学短视频的", LIGHT_GRAY, False),
    ("左下角链接", WHITE, True),
], sz=32, a=PP_ALIGN.CENTER, sp=1.8)

# ═══════════════════════════════════════════════════
# S13: 先这样
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 5.0, 8, 2.0, "先这样", 44, WHITE, True)
txt(s, 0.5, 7.0, 8, 1.2, "发了视频，我就正式开始", 28, LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# S14: 尾帧 Day 1 / 22
# ═══════════════════════════════════════════════════
s = dark()
txt(s, 0.5, 3.5, 8, 1.5, "Day 1 / 22", 56, WHITE, True)
line(s, 2, 5.0, 5, MID_GRAY, 2)
multi(s, 1.5, 5.5, 6, 4, [
    ("当前报名：0 人", LIGHT_GRAY, False),
    ("目标：1000 人", LIGHT_GRAY, False),
    ("距首场：14 天", LIGHT_GRAY, False),
], sz=24, a=PP_ALIGN.CENTER, sp=2.0)
line(s, 2, 9.0, 5, MID_GRAY, 1)
multi(s, 1, 9.5, 7, 3, [
    ("拓圈 Topoo", MID_GRAY, False),
    ("让小而美被世界看见", DIM_WHITE, False),
], sz=18, a=PP_ALIGN.CENTER, sp=1.6)

# ── 保存 ──
output = r"C:\Users\Administrator\WorkBuddy\Claw\temp_event\Day1-PPT素材-竖版9x16.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
