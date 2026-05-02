# -*- coding: utf-8 -*-
"""
句容市OPC大赛方案 - 政府汇报PPT
风格：黑白灰+红 高端商务风
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 配色方案：黑白灰 + 红色 ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT   = RGBColor(0xF7, 0xF7, 0xF7)
BG_CARD    = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT     = RGBColor(0xC8, 0x2A, 0x2A)   # 中国红
ACCENT_LIGHT = RGBColor(0xE0, 0x5A, 0x5A)
ACCENT_BG  = RGBColor(0xFD, 0xF0, 0xF0)   # 极浅红底
GRAY_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_MID   = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xBB, 0xBB, 0xBB)
GRAY_LINE  = RGBColor(0xE0, 0xE0, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)

FONT_CN = '微软雅黑'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 工具函数 ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color, line_color, line_width)

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color, line_color, line_width)

def add_circle(slide, left, top, size, fill_color, line_color=None):
    return add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color, line_color)

def set_shape_shadow(shape, blur=Pt(8), offset_x=Pt(2), offset_y=Pt(2), alpha=30000):
    """给shape添加柔和阴影"""
    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur)),
        'dist': str(int(Pt(3))),
        'dir': '5400000',
        'algn': 'tl',
        'rotWithShape': '0'
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alphaElem = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha)})
    srgbClr.append(alphaElem)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)

def add_text_box(slide, left, top, width, height, text, font_size=14, font_color=GRAY_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.5, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=14, font_color=GRAY_DARK, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.8, bullet=False, bullet_char="●"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet and line.strip():
            p.text = f"{bullet_char} {line}"
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_CN
        p.alignment = alignment
        p.space_after = Pt(3)
        p.space_before = Pt(1)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_divider(slide, left, top, width, color=ACCENT, thickness=Pt(2)):
    return add_rect(slide, left, top, width, thickness, color)

def add_page_header(slide, section_num, section_title):
    """统一的页头样式"""
    add_rect(slide, Inches(0), Inches(0), W, Pt(3), ACCENT)
    # 左侧留白区域 + 编号
    add_rect(slide, Inches(0), Inches(0), Inches(0.6), H, BG_DARK)
    # 编号竖排
    add_text_box(slide, Inches(0.08), Inches(0.4), Inches(0.44), Inches(0.8),
                 section_num, font_size=28, font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    # 分割竖线
    add_rect(slide, Inches(0.6), Inches(0), Pt(2), H, GRAY_LINE)
    # 标题
    add_text_box(slide, Inches(1.0), Inches(0.45), Inches(10), Inches(0.7),
                 section_title, font_size=30, font_color=GRAY_DARK, bold=True)
    add_divider(slide, Inches(1.0), Inches(1.25), Inches(2.2), ACCENT, Pt(3))
    # 页脚
    add_text_box(slide, Inches(1.0), Inches(6.95), Inches(5), Inches(0.3),
                 "句容市OPC大赛方案", font_size=9, font_color=GRAY_LIGHT)
    add_rect(slide, Inches(0.6), Inches(6.85), Inches(12.7), Pt(1), GRAY_LINE)

def make_stat_card(slide, left, top, width, height, num, label, color=ACCENT):
    """高端数据卡片"""
    card = add_rounded_rect(slide, left, top, width, height, WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, left, top, width, Pt(4), color)
    add_text_box(slide, left, top + Inches(0.25), width, Inches(0.7),
                 num, font_size=34, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.95), width, Inches(0.35),
                 label, font_size=13, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)

def make_content_card(slide, left, top, width, height, title, items, accent_color=ACCENT, title_size=18, item_size=13):
    """高端内容卡片"""
    card = add_rounded_rect(slide, left, top, width, height, WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    # 顶部色条
    add_rect(slide, left, top, width, Pt(4), accent_color)
    # 标题
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.5), Inches(0.45),
                 title, font_size=title_size, font_color=GRAY_DARK, bold=True)
    # 分割线
    add_divider(slide, left + Inches(0.3), top + Inches(0.7), Inches(0.8), accent_color, Pt(2))
    # 内容
    if items:
        add_multiline_box(slide, left + Inches(0.3), top + Inches(0.85), width - Inches(0.5), height - Inches(1.1),
                          items, font_size=item_size, font_color=GRAY_MID, bullet=True, line_spacing=1.8, bullet_char="▸")
    return card


# ═══════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# 左侧红色竖条
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

# 装饰横线
add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.3), Pt(1), GRAY_LIGHT)

# 主标题
add_text_box(slide, Inches(1.5), Inches(1.9), Inches(10.3), Inches(1.4),
             "句容市OPC大赛", font_size=52, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 英文副标题
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.5),
             "JURONG OPC COMPETITION", font_size=14, font_color=GRAY_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# 红色分割线
add_rect(slide, Inches(5.5), Inches(3.85), Inches(2.3), Pt(3), ACCENT)

# 口号
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
             "以赛招商  ·  以赛聚人  ·  以赛造势", font_size=22, font_color=ACCENT_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

# 装饰横线
add_rect(slide, Inches(1.5), Inches(5.0), Inches(10.3), Pt(1), GRAY_LIGHT)

# 组织信息 - 紧凑排列
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.35),
             "主办单位：句容市商务局  ·  句容市科技局", font_size=13, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.65), Inches(10.3), Inches(0.35),
             "承办单位：振邺科技  ·  智科数娱", font_size=13, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.35),
             "协办单位：香宁研究院  ·  龙虾联盟", font_size=13, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# 底部标注
add_text_box(slide, Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.3),
             "工作汇报方案（第一稿）  |  2025年度", font_size=10, font_color=RGBColor(0x55,0x55,0x55), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 第2页：目录
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_WHITE)
add_page_header(slide, "00", "目录")

toc_items = [
    ("01", "项目背景与意义"),
    ("02", "大赛概况"),
    ("03", "参赛对象与赛道"),
    ("04", "创作主题"),
    ("05", "赛程安排"),
    ("06", "评审体系"),
    ("07", "奖项设置"),
    ("08", "赋能与培训"),
    ("09", "宣传推广"),
    ("10", "预期效果"),
    ("11", "后续规划"),
]

# 左栏
for i, (num, title) in enumerate(toc_items[:6]):
    y = Inches(1.7) + Inches(i * 0.78)
    # 编号
    add_text_box(slide, Inches(1.2), y, Inches(0.6), Inches(0.45),
                 num, font_size=22, font_color=ACCENT if i == 0 else GRAY_LIGHT, bold=True)
    # 标题
    add_text_box(slide, Inches(1.9), y + Pt(4), Inches(4), Inches(0.4),
                 title, font_size=16, font_color=GRAY_DARK, bold=False)
    # 底线
    if i < 5:
        add_rect(slide, Inches(1.2), y + Inches(0.58), Inches(4.8), Pt(1), GRAY_LINE)

# 右栏
for i, (num, title) in enumerate(toc_items[6:]):
    y = Inches(1.7) + Inches(i * 0.78)
    add_text_box(slide, Inches(7.0), y, Inches(0.6), Inches(0.45),
                 num, font_size=22, font_color=GRAY_LIGHT, bold=True)
    add_text_box(slide, Inches(7.7), y + Pt(4), Inches(4), Inches(0.4),
                 title, font_size=16, font_color=GRAY_DARK, bold=False)
    if i < 4:
        add_rect(slide, Inches(7.0), y + Inches(0.58), Inches(4.8), Pt(1), GRAY_LINE)

# 右侧装饰区
add_rect(slide, Inches(6.5), Inches(1.5), Pt(1), Inches(5.2), GRAY_LINE)


# ═══════════════════════════════════════════
# 第3页：项目背景与意义
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "01", "项目背景与意义")

# 核心描述 - 引用块风格
quote_bg = add_rounded_rect(slide, Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.2), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(quote_bg)
add_rect(slide, Inches(1.0), Inches(1.6), Pt(5), Inches(1.2), ACCENT)
add_text_box(slide, Inches(1.4), Inches(1.75), Inches(10.5), Inches(0.9),
             '在"内容+电商"深度融合的当下，OPC（一人公司/个人创作者经济）已成为电商内容传播和产品推广的核心力量。通过举办OPC大赛，以赛事为抓手，批量吸引优质创作者关注句容、汇聚句容，为句容跨境电商产业园的招商运营注入持续的内容活力和人才储备。',
             font_size=14, font_color=GRAY_MID, line_spacing=1.9)

# 四大意义
cards_data = [
    ("城市名片", ["打造句容城市新名片", "展示产业活力与创新形象"], ACCENT),
    ("人才汇聚", ["批量吸引优质内容创作者", "建立OPC人才储备库"], GRAY_DARK),
    ("产业出圈", ["推动本地产品走向全国", "跨境电商产业园招商引流"], ACCENT),
    ("生态建设", ["构建内容电商生态链", "以赛促产、产城融合"], GRAY_MID),
]

card_w = Inches(2.7)
for i, (title, items, color) in enumerate(cards_data):
    x = Inches(1.0) + (card_w + Inches(0.25)) * i
    y = Inches(3.2)
    make_content_card(slide, x, y, card_w, Inches(2.8), title, items, accent_color=color, title_size=20, item_size=13)


# ═══════════════════════════════════════════
# 第4页：大赛概况
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "02", "大赛概况")

# 信息卡片 - 2x2网格
info_data = [
    ("大赛名称", "句容市OPC大赛（名称待定）", "01", ACCENT),
    ("参赛人群", "高校赛道 + 社会赛道 · 双轨并行", "02", GRAY_DARK),
    ("比赛形式", "多赛道内容创作竞赛\n覆盖OPC核心能力方向", "03", ACCENT),
    ("赛事定位", "以赛招商、以赛聚人、以赛造势\n为跨境电商产业园注入活力", "04", GRAY_MID),
]

for i, (title, desc, num, color) in enumerate(info_data):
    col = i % 2
    row = i // 2
    x = Inches(1.0) + (Inches(5.8) + Inches(0.4)) * col
    y = Inches(1.6) + (Inches(2.2) + Inches(0.25)) * row
    card = add_rounded_rect(slide, x, y, Inches(5.8), Inches(2.2), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    # 编号
    add_rounded_rect(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.28), Inches(0.55), Inches(0.5),
                 num, font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, x + Inches(1.0), y + Inches(0.3), Inches(4.5), Inches(0.45),
                 title, font_size=20, font_color=GRAY_DARK, bold=True)
    # 描述
    add_text_box(slide, x + Inches(0.25), y + Inches(1.0), Inches(5.3), Inches(1.0),
                 desc, font_size=14, font_color=GRAY_MID, line_spacing=1.8)

# 组织架构
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(3), Inches(0.4),
             "组织架构", font_size=16, font_color=GRAY_DARK, bold=True)

org_data = [
    ("主办单位", "句容市商务局 · 句容市科技局", ACCENT),
    ("承办单位", "振邺科技 · 智科数娱", GRAY_DARK),
    ("协办单位", "香宁研究院 · 龙虾联盟", GRAY_MID),
]
for i, (role, units, color) in enumerate(org_data):
    x = Inches(5.5) + Inches(i * 2.7)
    add_rect(slide, x, Inches(6.35), Pt(4), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.2), Inches(6.35), Inches(2.4), Inches(0.35),
                 f"{role}：{units}", font_size=11, font_color=GRAY_MID)


# ═══════════════════════════════════════════
# 第5页：参赛对象与赛道
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "03", "参赛对象与赛道")

# 左侧：参赛对象
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(4), Inches(0.4),
             "参赛对象", font_size=20, font_color=GRAY_DARK, bold=True)
add_divider(slide, Inches(1.0), Inches(2.05), Inches(0.8), ACCENT, Pt(2))

# 高校赛道
card1 = add_rounded_rect(slide, Inches(1.0), Inches(2.4), Inches(5.5), Inches(2.0), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card1)
add_rect(slide, Inches(1.0), Inches(2.4), Inches(5.5), Pt(4), ACCENT)
add_text_box(slide, Inches(1.3), Inches(2.65), Inches(5), Inches(0.4),
             "高校赛道", font_size=18, font_color=ACCENT, bold=True)
add_multiline_box(slide, Inches(1.3), Inches(3.15), Inches(4.8), Inches(1.0),
                  ["面向全国高校在校生（重点覆盖江苏及长三角）",
                   "个人或团队参赛，每队不超过3人",
                   "高校组织推荐报名"],
                  font_size=12, font_color=GRAY_MID, bullet=True, line_spacing=1.7, bullet_char="▸")

# 社会赛道
card2 = add_rounded_rect(slide, Inches(1.0), Inches(4.7), Inches(5.5), Inches(2.0), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card2)
add_rect(slide, Inches(1.0), Inches(4.7), Inches(5.5), Pt(4), GRAY_DARK)
add_text_box(slide, Inches(1.3), Inches(4.95), Inches(5), Inches(0.4),
             "社会赛道", font_size=18, font_color=GRAY_DARK, bold=True)
add_multiline_box(slide, Inches(1.3), Inches(5.45), Inches(4.8), Inches(1.0),
                  ["面向社会内容创作者、自媒体从业者",
                   "鼓励有意向从事内容电商的个人和团队",
                   "句容本地创业者优先"],
                  font_size=12, font_color=GRAY_MID, bullet=True, line_spacing=1.7, bullet_char="▸")

# 右侧：四大赛道
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "四大赛道", font_size=20, font_color=GRAY_DARK, bold=True)
add_divider(slide, Inches(7.2), Inches(2.05), Inches(0.8), ACCENT, Pt(2))

tracks = [
    ("01", "短视频创作", "种草 / 测评 / 带货 / 品牌宣传", ACCENT),
    ("02", "短剧创作", "产品故事 / 品牌理念 / 城市文化", GRAY_DARK),
    ("03", "直播实战", "直播带货 / 选品话术 / 互动转化", ACCENT_LIGHT),
    ("04", "内容电商", "图文种草 / 多形式产品推广", GRAY_MID),
]

for i, (num, name, desc, color) in enumerate(tracks):
    y = Inches(2.4) + Inches(i * 1.1)
    card = add_rounded_rect(slide, Inches(7.2), y, Inches(5.3), Inches(0.9), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    # 编号圆
    add_circle(slide, Inches(7.5), y + Inches(0.17), Inches(0.55), color)
    add_text_box(slide, Inches(7.5), y + Inches(0.2), Inches(0.55), Inches(0.5),
                 num, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 赛道名
    add_text_box(slide, Inches(8.3), y + Inches(0.1), Inches(3.8), Inches(0.4),
                 name, font_size=17, font_color=GRAY_DARK, bold=True)
    add_text_box(slide, Inches(8.3), y + Inches(0.5), Inches(3.8), Inches(0.3),
                 desc, font_size=12, font_color=GRAY_MID)


# ═══════════════════════════════════════════
# 第6页：创作主题
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "04", "创作主题")

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.4),
             "大赛主题围绕产品推广和城市文旅两大方向，具体选题结合赞助商产品和本地产业资源灵活设置",
             font_size=13, font_color=GRAY_MID)

# 两大方向 - 对称大卡片
# 左
card1 = add_rounded_rect(slide, Inches(1.0), Inches(2.3), Inches(5.6), Inches(4.0), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card1)
add_rect(slide, Inches(1.0), Inches(2.3), Inches(5.6), Inches(0.06), ACCENT)
# 标签
add_rounded_rect(slide, Inches(1.4), Inches(2.65), Inches(1.2), Inches(0.4), ACCENT)
add_text_box(slide, Inches(1.4), Inches(2.67), Inches(1.2), Inches(0.35),
             "方向一", font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), Inches(2.65), Inches(3.5), Inches(0.4),
             "好物推荐", font_size=24, font_color=GRAY_DARK, bold=True)
add_text_box(slide, Inches(1.4), Inches(3.25), Inches(4.8), Inches(0.35),
             "围绕大赛指定产品 / 赞助商产品进行内容创作", font_size=13, font_color=GRAY_MID)
add_multiline_box(slide, Inches(1.4), Inches(3.8), Inches(4.8), Inches(2.2),
                  ["内容种草：产品亮点挖掘与创意呈现",
                   "产品测评：使用体验深度展示",
                   "场景展示：真实使用场景内容创作",
                   "品牌故事：产品背后的品牌理念传达"],
                  font_size=13, font_color=GRAY_MID, bullet=True, line_spacing=2.0, bullet_char="▸")

# 右
card2 = add_rounded_rect(slide, Inches(6.9), Inches(2.3), Inches(5.6), Inches(4.0), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card2)
add_rect(slide, Inches(6.9), Inches(2.3), Inches(5.6), Inches(0.06), GRAY_DARK)
add_rounded_rect(slide, Inches(7.3), Inches(2.65), Inches(1.2), Inches(0.4), GRAY_DARK)
add_text_box(slide, Inches(7.3), Inches(2.67), Inches(1.2), Inches(0.35),
             "方向二", font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.7), Inches(2.65), Inches(3.5), Inches(0.4),
             "城市文旅", font_size=24, font_color=GRAY_DARK, bold=True)
add_text_box(slide, Inches(7.3), Inches(3.25), Inches(4.8), Inches(0.35),
             "围绕句容城市风貌、特色产业、文化故事", font_size=13, font_color=GRAY_MID)
add_multiline_box(slide, Inches(7.3), Inches(3.8), Inches(4.8), Inches(2.2),
                  ["城市风貌：句容自然风光与城市景观",
                   "特色产业：本地优势产业深度展示",
                   "文化故事：历史人文与非遗传承",
                   "城市活力：现代发展与青年视角"],
                  font_size=13, font_color=GRAY_MID, bullet=True, line_spacing=2.0, bullet_char="▸")

# 底部提示
add_rounded_rect(slide, Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.5), ACCENT_BG, line_color=ACCENT, line_width=Pt(1))
add_text_box(slide, Inches(1.5), Inches(6.55), Inches(10.5), Inches(0.35),
             "*  具体创作选题将在赛事启动时公布，结合赞助商产品和本地产业资源灵活设置",
             font_size=12, font_color=ACCENT)


# ═══════════════════════════════════════════
# 第7页：赛程安排（5月-12月）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "05", "赛程安排")

# 总时间轴标注
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11.5), Inches(0.35),
             "2025年5月  —  2025年12月   |   全程8个月", font_size=14, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)

# 主时间线
phases = [
    ("5月", "新闻发布\n暨启动仪式", "举办新闻发布会\n政府领导致辞\n赛事正式启动\n全渠道宣发开启", ACCENT),
    ("6-7月", "选拔赛", "线上内容创作\n高校/社会双轨评选\n各赛道选出晋级选手\n阶段性数据公示", GRAY_DARK),
    ("9-10月", "晋级赛", "限时主题创作\n赛前集训营\n导师一对一指导\n选出决赛选手", ACCENT_LIGHT),
    ("12月", "总决赛\n暨颁奖典礼", "现场创作展示\n直播实战环节\n颁奖典礼\n年度OPC盛典", ACCENT),
]

# 时间线横轴
timeline_y = Inches(2.3)
add_rect(slide, Inches(1.5), timeline_y + Inches(0.22), Inches(10.5), Pt(3), GRAY_LINE)

phase_w = Inches(2.5)
for i, (time, name, desc, color) in enumerate(phases):
    x = Inches(0.8) + (phase_w + Inches(0.3)) * i
    
    # 时间节点
    add_circle(slide, x + Inches(1.05), timeline_y + Inches(0.05), Inches(0.38), color)
    add_text_box(slide, x + Inches(1.05), timeline_y + Inches(0.08), Inches(0.38), Inches(0.32),
                 str(i+1), font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 月份
    add_text_box(slide, x, timeline_y - Inches(0.35), phase_w, Inches(0.3),
                 time, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 内容卡片
    card = add_rounded_rect(slide, x, Inches(2.95), phase_w, Inches(3.6), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, x, Inches(2.95), phase_w, Pt(4), color)
    # 阶段名
    add_text_box(slide, x + Inches(0.2), Inches(3.15), phase_w - Inches(0.4), Inches(0.7),
                 name, font_size=17, font_color=GRAY_DARK, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.4)
    # 分割线
    add_divider(slide, x + Inches(0.3), Inches(3.9), Inches(1.2), color, Pt(2))
    # 描述
    add_multiline_box(slide, x + Inches(0.2), Inches(4.1), phase_w - Inches(0.4), Inches(2.2),
                      desc.split('\n'), font_size=12, font_color=GRAY_MID, bullet=True, line_spacing=1.9, bullet_char="▸")


# ═══════════════════════════════════════════
# 第8页：评审体系
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "06", "评审体系")

# 五维评审 - 环形进度条风格
criteria = [
    ("内容创意", 25, "创意独特性、故事性、可看性", ACCENT),
    ("制作质量", 20, "画面质量、剪辑节奏、表达力", GRAY_DARK),
    ("数据表现", 20, "播放量、互动率、完播率/转化率", ACCENT_LIGHT),
    ("商业转化潜力", 20, "选品能力、变现可能性", GRAY_MID),
    ("主题契合度", 15, "与指定主题的契合度", GRAY_LIGHT),
]

bar_start_x = Inches(3.5)
bar_max_w = Inches(5.5)

for i, (name, weight, desc, color) in enumerate(criteria):
    y = Inches(1.8) + Inches(i * 0.95)
    # 名称
    add_text_box(slide, Inches(1.0), y, Inches(2.2), Inches(0.4),
                 name, font_size=17, font_color=GRAY_DARK, bold=True)
    # 进度条背景
    add_rounded_rect(slide, bar_start_x, y + Inches(0.06), bar_max_w, Inches(0.28), GRAY_LINE)
    # 进度条
    bar_w = int(bar_max_w * weight / 25)
    bar = add_rounded_rect(slide, bar_start_x, y + Inches(0.06), bar_w, Inches(0.28), color)
    # 权重
    add_text_box(slide, bar_start_x + bar_max_w + Inches(0.3), y, Inches(0.8), Inches(0.4),
                 f"{weight}%", font_size=18, font_color=color, bold=True)
    # 描述
    add_text_box(slide, Inches(1.0), y + Inches(0.38), Inches(2.2), Inches(0.3),
                 desc, font_size=10, font_color=GRAY_MID)

# 评审团
add_rect(slide, Inches(1.0), Inches(6.1), Inches(11.5), Inches(0.6), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.5), Inches(0.4),
             "评审团构成：行业专家  ·  平台代表  ·  品牌方代表  ·  政府代表",
             font_size=14, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 第9页：奖项设置
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "07", "奖项设置")

# 高校赛道
card1 = add_rounded_rect(slide, Inches(1.0), Inches(1.6), Inches(5.6), Inches(3.8), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card1)
add_rect(slide, Inches(1.0), Inches(1.6), Inches(5.6), Inches(0.06), ACCENT)
add_text_box(slide, Inches(1.5), Inches(1.9), Inches(4.5), Inches(0.4),
             "高校赛道", font_size=20, font_color=GRAY_DARK, bold=True)

uni_prizes = [
    ("一等奖", "各赛道1名", "奖金 + 证书 + 入园扶持", ACCENT),
    ("二等奖", "各赛道3名", "奖金 + 证书 + 实习机会", GRAY_DARK),
    ("三等奖", "各赛道5名", "奖金 + 证书", GRAY_MID),
    ("优秀奖", "若干", "荣誉证书", GRAY_LIGHT),
]

for i, (prize, quota, reward, color) in enumerate(uni_prizes):
    y = Inches(2.5) + Inches(i * 0.7)
    # 色块标识
    add_rect(slide, Inches(1.5), y + Inches(0.05), Pt(4), Inches(0.4), color)
    add_text_box(slide, Inches(1.8), y, Inches(1.5), Inches(0.35),
                 prize, font_size=16, font_color=GRAY_DARK, bold=True)
    add_text_box(slide, Inches(3.3), y, Inches(1.5), Inches(0.35),
                 quota, font_size=13, font_color=GRAY_MID)
    add_text_box(slide, Inches(1.8), y + Inches(0.32), Inches(4.5), Inches(0.3),
                 reward, font_size=11, font_color=GRAY_MID)

# 社会赛道
card2 = add_rounded_rect(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(3.8), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
set_shape_shadow(card2)
add_rect(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(0.06), GRAY_DARK)
add_text_box(slide, Inches(7.4), Inches(1.9), Inches(4.5), Inches(0.4),
             "社会赛道", font_size=20, font_color=GRAY_DARK, bold=True)

soc_prizes = [
    ("一等奖", "各赛道1名", "奖金 + 证书 + 商务合作机会", ACCENT),
    ("二等奖", "各赛道3名", "奖金 + 证书 + 流量扶持", GRAY_DARK),
    ("三等奖", "各赛道5名", "奖金 + 证书", GRAY_MID),
]

for i, (prize, quota, reward, color) in enumerate(soc_prizes):
    y = Inches(2.5) + Inches(i * 0.7)
    add_rect(slide, Inches(7.4), y + Inches(0.05), Pt(4), Inches(0.4), color)
    add_text_box(slide, Inches(7.7), y, Inches(1.5), Inches(0.35),
                 prize, font_size=16, font_color=GRAY_DARK, bold=True)
    add_text_box(slide, Inches(9.2), y, Inches(1.5), Inches(0.35),
                 quota, font_size=13, font_color=GRAY_MID)
    add_text_box(slide, Inches(7.7), y + Inches(0.32), Inches(4.5), Inches(0.3),
                 reward, font_size=11, font_color=GRAY_MID)

# 特别奖项
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(3), Inches(0.4),
             "特别奖项", font_size=18, font_color=GRAY_DARK, bold=True)

special = [("最佳创意奖", ACCENT), ("最佳人气奖", ACCENT_LIGHT), ("最佳商业潜力奖", GRAY_DARK)]
for i, (name, color) in enumerate(special):
    x = Inches(1.0) + Inches(i * 4.0)
    card = add_rounded_rect(slide, x, Inches(6.2), Inches(3.6), Inches(0.55), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, x, Inches(6.2), Pt(5), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.3), Inches(6.27), Inches(3), Inches(0.4),
                 name, font_size=15, font_color=GRAY_DARK, bold=True)


# ═══════════════════════════════════════════
# 第10页：赋能与培训
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "08", "赋能与培训")

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.35),
             "全链条赋能体系，从赛前启蒙到赛后持续成长",
             font_size=13, font_color=GRAY_MID)

train_data = [
    ("赛前", "基础赋能", [
        "高校巡回宣讲会",
        "OPC行业趋势与新职业前景",
        "线上培训课程",
        "内容创作基础入门"
    ], ACCENT),
    ("赛中", "进阶提升", [
        "晋级赛前集训营",
        "进阶创作技巧",
        "选品逻辑与数据运营",
        "行业导师一对一指导"
    ], GRAY_DARK),
    ("赛后", "持续成长", [
        "获奖选手纳入人才库",
        "对接跨境电商产业园",
        "优秀作品持续传播",
        "为产业园提供内容资产"
    ], ACCENT_LIGHT),
]

card_w = Inches(3.7)
for i, (phase, sub_title, items, color) in enumerate(train_data):
    x = Inches(1.0) + (card_w + Inches(0.3)) * i
    y = Inches(2.2)
    card = add_rounded_rect(slide, x, y, card_w, Inches(4.2), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, x, y, card_w, Pt(4), color)
    # 阶段标签
    add_rounded_rect(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.85), Inches(0.38), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.27), Inches(0.85), Inches(0.33),
                 phase, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(1.3), y + Inches(0.27), Inches(2.2), Inches(0.38),
                 sub_title, font_size=17, font_color=GRAY_DARK, bold=True)
    # 分割线
    add_divider(slide, x + Inches(0.25), y + Inches(0.8), Inches(0.8), color, Pt(2))
    # 内容
    add_multiline_box(slide, x + Inches(0.25), y + Inches(1.0), card_w - Inches(0.5), Inches(2.8),
                      items, font_size=13, font_color=GRAY_MID, bullet=True, line_spacing=2.0, bullet_char="▸")

# 连接箭头
for i in range(2):
    x = Inches(4.7) + Inches(i * 4.0)
    add_text_box(slide, x, Inches(3.9), Inches(0.5), Inches(0.5),
                 "▸", font_size=28, font_color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 第11页：宣传推广
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "09", "宣传推广")

promo_data = [
    ("线上推广", "多平台流量矩阵", [
        "抖音话题挑战赛",
        "高校自媒体矩阵联动",
        "微信/视频号/小红书分发",
        "KOL/KOC联合推广"
    ], ACCENT),
    ("线下推广", "深度触达目标人群", [
        "高校校园宣传",
        "句容本地商圈社区宣传",
        "新闻发布会暨启动仪式",
        "高铁站/公交站广告投放"
    ], GRAY_DARK),
    ("媒体传播", "权威背书与行业影响", [
        "省/市级主流媒体报道",
        "行业媒体专题报道",
        "高校媒体联动传播",
        "行业自媒体KOL发声"
    ], ACCENT_LIGHT),
]

card_w = Inches(3.7)
for i, (title, sub, items, color) in enumerate(promo_data):
    x = Inches(1.0) + (card_w + Inches(0.3)) * i
    y = Inches(1.7)
    card = add_rounded_rect(slide, x, y, card_w, Inches(4.7), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, x, y, card_w, Pt(4), color)
    # 编号
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.4),
                 f"0{i+1}", font_size=14, font_color=color, bold=True)
    # 标题
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.1), Inches(0.45),
                 title, font_size=20, font_color=GRAY_DARK, bold=True)
    # 副标题
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.3),
                 sub, font_size=12, font_color=color)
    # 分割线
    add_divider(slide, x + Inches(0.3), y + Inches(1.55), Inches(0.8), color, Pt(2))
    # 内容
    add_multiline_box(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.1), Inches(2.5),
                      items, font_size=13, font_color=GRAY_MID, bullet=True, line_spacing=2.0, bullet_char="▸")


# ═══════════════════════════════════════════
# 第12页：预期效果
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "10", "预期效果")

# 数据大屏
stats = [
    ("2,500+", "总参赛人数", ACCENT),
    ("30+", "覆盖高校数量", GRAY_DARK),
    ("5,000万+", "话题播放量", ACCENT_LIGHT),
    ("1亿+", "全网曝光量", ACCENT),
    ("5,000+", "优质内容产出", GRAY_MID),
]

for i, (num, label, color) in enumerate(stats):
    x = Inches(0.7) + Inches(i * 2.5)
    make_stat_card(slide, x, Inches(1.7), Inches(2.2), Inches(1.5), num, label, color)

# 产业价值
add_text_box(slide, Inches(1.0), Inches(3.7), Inches(3), Inches(0.4),
             "产业价值", font_size=18, font_color=GRAY_DARK, bold=True)
add_divider(slide, Inches(1.0), Inches(4.1), Inches(0.8), ACCENT, Pt(2))

value_items = [
    ("人才储备", "为跨境电商产业园\n建立OPC人才库", ACCENT),
    ("产业对接", "获奖创作者优先\n对接园区企业", GRAY_DARK),
    ("品牌IP", "赛事形成年度IP\n持续输送人才和内容", ACCENT_LIGHT),
]

for i, (title, desc, color) in enumerate(value_items):
    x = Inches(1.0) + Inches(i * 4.0)
    y = Inches(4.4)
    card = add_rounded_rect(slide, x, y, Inches(3.7), Inches(2.0), WHITE, line_color=GRAY_LINE, line_width=Pt(1))
    set_shape_shadow(card)
    add_rect(slide, x, y, Inches(3.7), Pt(4), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(3.1), Inches(0.4),
                 title, font_size=18, font_color=GRAY_DARK, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.1), Inches(1.0),
                 desc, font_size=14, font_color=GRAY_MID, line_spacing=1.8)


# ═══════════════════════════════════════════
# 第13页：后续规划
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_page_header(slide, "11", "后续规划")

plans = [
    ("01", "创作者入驻", "获奖创作者对接跨境电商产业园，成为首批内容入驻企业/个人", ACCENT),
    ("02", "人才池运营", "参赛选手社群持续运营，定期举办活动，保持人才池活跃", GRAY_DARK),
    ("03", "年度IP化", "大赛年度IP化，每年举办一届，形成持续品牌影响力", ACCENT_LIGHT),
    ("04", "模式复制", "模式跑通后可复制到其他产业园区，实现规模化输出", GRAY_MID),
]

for i, (num, title, desc, color) in enumerate(plans):
    y = Inches(1.8) + Inches(i * 1.2)
    # 左侧编号
    card = add_rounded_rect(slide, Inches(1.0), y, Inches(0.75), Inches(0.75), color)
    set_shape_shadow(card)
    add_text_box(slide, Inches(1.0), y + Pt(8), Inches(0.75), Inches(0.55),
                 num, font_size=22, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, Inches(2.1), y + Pt(2), Inches(4), Inches(0.5),
                 title, font_size=22, font_color=GRAY_DARK, bold=True)
    # 描述
    add_text_box(slide, Inches(2.1), y + Inches(0.45), Inches(9.5), Inches(0.5),
                 desc, font_size=14, font_color=GRAY_MID)
    # 连接线
    if i < 3:
        add_rect(slide, Inches(1.35), y + Inches(0.75), Pt(2), Inches(0.45), GRAY_LINE)

# 底部愿景
add_rounded_rect(slide, Inches(1.0), Inches(6.2), Inches(11.5), Inches(0.65), BG_DARK)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.5), Inches(0.45),
             "以赛为始，以产为终  ——  打造句容内容电商新生态",
             font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 第14页：封底
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# 左侧红色竖条
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

# 装饰线
add_rect(slide, Inches(1.5), Inches(2.0), Inches(10.3), Pt(1), GRAY_LIGHT)

add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.9),
             "感谢聆听", font_size=48, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_divider(slide, Inches(5.8), Inches(3.4), Inches(1.7), ACCENT, Pt(3))

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.4),
             "句容市OPC大赛", font_size=18, font_color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.4),
             "以赛招商  ·  以赛聚人  ·  以赛造势", font_size=14, font_color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# 装饰线
add_rect(slide, Inches(1.5), Inches(4.8), Inches(10.3), Pt(1), GRAY_LIGHT)

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.3),
             "主办单位：句容市商务局  ·  句容市科技局", font_size=12, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.3),
             "承办单位：振邺科技  ·  智科数娱", font_size=12, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.3),
             "协办单位：香宁研究院  ·  龙虾联盟", font_size=12, font_color=GRAY_MID, alignment=PP_ALIGN.CENTER)


# ── 输出 ──
output_path = r'c:\Users\Administrator\WorkBuddy\Claw\句容市OPC大赛方案.pptx'
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
