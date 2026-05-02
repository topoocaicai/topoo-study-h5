# -*- coding: utf-8 -*-
"""
Day 1 视频PPT —— "我要在南京8天拉1000人"
黑白灰高级风格，用于视频录制的PPT翻页素材
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 颜色系统 ───
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_BG = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT_RED = RGBColor(0xFF, 0x44, 0x44)
DIM_WHITE = RGBColor(0x99, 0x99, 0x99)

SLIDE_W = Inches(16)
SLIDE_H = Inches(9)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── 工具函数 ───
def add_black_slide(prs):
    """添加纯黑背景幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide

def add_white_slide(prs):
    """添加白底幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=40,
                 color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
                 font_name="Microsoft YaHei", line_spacing=1.5):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multiline_text(slide, left, top, width, height, lines,
                       font_size=32, color=WHITE, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei",
                       line_spacing=1.8):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 支持 (text, color, bold) 元组
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1] if len(line) > 1 else color
            p.font.bold = line[2] if len(line) > 2 else bold
        else:
            p.text = line
            p.font.color.rgb = color
            p.font.bold = bold
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        if line_spacing:
            p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_line(slide, left, top, width, color=MID_GRAY, thickness=1):
    """添加分割线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(thickness)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════
# Slide 1: 钩子 —— 南京 / 8天 / 1000人
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 0, 1.5, 16, 1.5, "南京", 72, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 3.0, 16, 1.5, "8天", 72, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 4.5, 16, 2.0, "1000人", 96, WHITE, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: 自我介绍
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 2, 3.0, 12, 1.2,
             "大家好，我是苏菜菜", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 2, 4.2, 12, 1.2,
             "拓圈Topoo发起者", 36, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 3: 打脸宣言
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 1, 3.2, 14, 2.0,
             "我要做一件大概率会被打脸的事", 44, WHITE, True, PP_ALIGN.CENTER)

# ── 穿插镜头①：手机录屏（此处放截图占位）──
# 实际录制时这里暂停，切手机录屏2秒

# ═══════════════════════════════════════════
# Slide 4: 活动信息
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_multiline_text(slide, 3, 1.5, 10, 6.5, [
    ("5月10日 — 5月17日", WHITE, True),
    ("", WHITE, False),
    ("南京 · 国创园", WHITE, False),
    ("6场线下培训", WHITE, False),
    ("", WHITE, False),
    ("讲师：张明军", LIGHT_GRAY, False),
    ("百万粉丝科技博主", MID_GRAY, False),
], font_size=36, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

# ── 穿插镜头②：张明军主页（此处暂停，切截图2秒）──

# ═══════════════════════════════════════════
# Slide 5: 南京940万 vs 1000人
# ═══════════════════════════════════════════
slide = add_white_slide(prs)
add_text_box(slide, 0, 2.0, 16, 1.2,
             "南京常住人口", 36, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0, 3.0, 16, 1.5,
             "940万", 80, RGBColor(0x33, 0x33, 0x33), True, PP_ALIGN.CENTER)
add_text_box(slide, 0, 5.0, 16, 1.2,
             "1000人 = 万分之一", 32, MID_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 6: "每一步都在掉人"
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 1, 0.8, 14, 1.2,
             "但做过线下活动的人都知道", 32, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_line(slide, 4, 2.0, 8, MID_GRAY, 1)
add_multiline_text(slide, 2.5, 2.5, 11, 5.5, [
    ("让一个人——", WHITE, False),
    ("看到信息", WHITE, False),
    ("↓", MID_GRAY, False),
    ("付 199 块", WHITE, False),
    ("↓", MID_GRAY, False),
    ("周末爬起来出门", WHITE, False),
    ("↓", MID_GRAY, False),
    ("坐地铁到场", WHITE, False),
    ("↓", MID_GRAY, False),
    ("坐 8 个小时", WHITE, False),
    ("", WHITE, False),
    ("每一步，都在掉人。", ACCENT_RED, True),
], font_size=30, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# ── 穿插镜头③：空场地实拍（此处暂停，切实拍3秒）──

# ═══════════════════════════════════════════
# Slide 7: 算账
# ═══════════════════════════════════════════
slide = add_white_slide(prs)
add_text_box(slide, 0, 2.5, 16, 2.0,
             "6 场 × 167 人 = 1002 人", 60,
             RGBColor(0x33, 0x33, 0x33), True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 8: 距首场 + 报名人数0
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 0, 1.5, 16, 1.0,
             "距首场：14天", 32, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0, 3.0, 16, 1.2,
             "当前报名人数：", 40, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0, 4.2, 16, 2.5,
             "0", 140, WHITE, True, PP_ALIGN.CENTER)

# ── 穿插镜头④：手部特写看后台（此处暂停，切实拍2秒）──

# ═══════════════════════════════════════════
# Slide 9: 每天更新真实数据
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 1, 2.5, 14, 1.5,
             "从今天开始", 40, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 14, 2.0,
             "我每天更新真实数据", 48, WHITE, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 10: 做到了算本事 / 没做到复盘素材
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_multiline_text(slide, 2, 2.0, 12, 5.5, [
    ("做到了 → 算本事", WHITE, True),
    ("", WHITE, False),
    ("没做到 → 复盘素材", LIGHT_GRAY, True),
    ("", WHITE, False),
    ("反正，不会无聊。", MID_GRAY, False),
], font_size=40, alignment=PP_ALIGN.CENTER, line_spacing=1.8)

# ═══════════════════════════════════════════
# Slide 11: 报名信息
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_multiline_text(slide, 3, 1.5, 10, 6.0, [
    ("南京 · 国创园", WHITE, True),
    ("", WHITE, False),
    ("5月10日 — 5月17日", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("¥199 / 人", WHITE, True),
    ("全天 6 大模块", LIGHT_GRAY, False),
    ("线下授课 · 张明军", MID_GRAY, False),
    ("", WHITE, False),
    ("点 链 接 ， 报 名 。", WHITE, True),
], font_size=32, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# ── 穿插镜头⑤：场地一角（此处暂停，切实拍2秒）──

# ═══════════════════════════════════════════
# Slide 12: 尾帧 —— Day 1 / 22
# ═══════════════════════════════════════════
slide = add_black_slide(prs)
add_text_box(slide, 0, 1.0, 16, 1.5,
             "Day 1 / 22", 72, WHITE, True, PP_ALIGN.CENTER)
add_line(slide, 4, 2.8, 8, MID_GRAY, 2)

add_multiline_text(slide, 3, 3.2, 10, 3.5, [
    ("当前报名：0 人", LIGHT_GRAY, False),
    ("目标：1000 人", LIGHT_GRAY, False),
    ("距首场：14 天", LIGHT_GRAY, False),
], font_size=30, alignment=PP_ALIGN.CENTER, line_spacing=2.0)

add_line(slide, 4, 6.2, 8, MID_GRAY, 1)
add_text_box(slide, 0, 6.8, 16, 1.0,
             "拓圈 Topoo", 28, MID_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 0, 7.5, 16, 1.0,
             "让小而美被世界看见", 22, DIM_WHITE, False, PP_ALIGN.CENTER)

# ── 保存 ──
output_path = r"C:\Users\Administrator\WorkBuddy\Claw\拓圈·1000人挑战·Day1-PPT素材.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
